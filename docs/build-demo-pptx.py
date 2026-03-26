"""Generate demo-questions.pptx for the NC Governor's Communications Agent."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# -- Colors --
NAVY = RGBColor(0x0A, 0x24, 0x63)
BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xD8, 0x3B, 0x01)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Reduce corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_question_slide(prs, title, subtitle, color, questions):
    """Add a capability question slide with a colored header bar and question list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Header bar
    bar = add_rounded_rect(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.9), color)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"

    # Subtitle
    add_text_box(slide, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
                 subtitle, font_size=14, color=GRAY, italic=True)

    # Questions
    y = Inches(1.95)
    row_height = Inches(0.82)

    for i, (prompt, why) in enumerate(questions):
        # Number circle
        circle_size = Inches(0.35)
        from pptx.enum.shapes import MSO_SHAPE
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.06), circle_size, circle_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = str(i + 1)
        cp.font.size = Pt(13)
        cp.font.color.rgb = WHITE
        cp.font.bold = True
        cp.font.name = "Segoe UI"
        cp.alignment = PP_ALIGN.CENTER

        # Question text
        txBox = slide.shapes.add_textbox(Inches(1.15), y, Inches(11.5), row_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = prompt
        p1.font.size = Pt(18)
        p1.font.color.rgb = DARK
        p1.font.bold = True
        p1.font.name = "Segoe UI"
        p1.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = why
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.font.name = "Segoe UI"

        y += row_height

    # Footer
    add_text_box(slide, Inches(0.4), Inches(7.0), Inches(12), Inches(0.4),
                 "NC DIT AI Office", font_size=9, color=RGBColor(0xBB, 0xBB, 0xBB))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ===== SLIDE 1: TITLE =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Demo Questions", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 "NC Governor's Communications Agent", font_size=26, color=RGBColor(0xBB, 0xCC, 0xFF),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
                 "Sample prompts for the live demo  |  March 2026", font_size=16,
                 color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)

    # ===== SLIDE 2: NEWS CLIPS =====
    clips_questions = [
        ("What clips came in this week about broadband?",
         "Starter from the talk track. Shows basic semantic search."),
        ("Show me the latest news.",
         "Triggers browse mode \u2014 returns most recent clips sorted by date."),
        ("Has the Governor said anything about rural internet access?",
         "Semantic match: finds \"broadband\" articles even though the words don't overlap."),
        ("Any press releases about economic development or jobs?",
         "Broad topic. Shows hybrid search combining meaning and keyword matching."),
        ("What's in the news about Hurricane Helene recovery?",
         "Timely topic. Matches western NC disaster recovery press releases."),
        ("Are there any clips about education or schools?",
         "Wide query \u2014 could match teacher pay, school safety, student wellbeing."),
    ]
    add_question_slide(prs, "News Clips", "Search Governor press releases scraped from governor.nc.gov", BLUE, clips_questions)

    # ===== SLIDE 3: REMARKS SEARCH =====
    remarks_questions = [
        ("What has the Governor said about mental health and public safety?",
         "Starter from the talk track. Hits the October column \u2014 mental health as public safety."),
        ("What language has the Governor used about keeping kids safe online?",
         "Direct hit on the September column. Shows retrieval of specific policy language."),
        ("How has the Governor talked about emergency preparedness and hurricane recovery?",
         "Pulls from August (prep) and November (Helene gratitude). Cross-speech synthesis."),
        ("What has the Governor said about law enforcement and supporting police?",
         "October column: $195M public safety package, raising officer salaries."),
        ("How does the Governor talk about North Carolina's history and democracy?",
         "January column: America 250, Halifax Resolves, Bill of Rights. Rich synthesis."),
        ("What has the Governor said about food insecurity and food banks?",
         "December column (MANNA, Second Harvest) + November (farmers). Quotes with citations."),
    ]
    add_question_slide(prs, "Remarks Search", "Search the Governor's speeches and columns \u2014 GPT-4o writes a synthesis with direct quotes", GREEN, remarks_questions)

    # ===== SLIDE 4: PROOFREAD =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = add_rounded_rect(slide, Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.9), ORANGE)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(20)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Transcript Proofreading"
    p.font.size = Pt(30)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"

    add_text_box(slide, Inches(0.6), Inches(1.35), Inches(12), Inches(0.5),
                 "Paste this sample transcript into the agent. The AI fixes errors and returns every change with a reason and confidence level.",
                 font_size=14, color=GRAY, italic=True)

    # Sample transcript box
    sample_rect = add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.8),
                                   RGBColor(0xFD, 0xF6, 0xEC))
    sample_rect.line.color.rgb = RGBColor(0xE8, 0xD5, 0xB0)
    sample_rect.line.width = Pt(1)

    sample_text = (
        'Good evening, North Carolinains. Im honored to adress this chamber tonite. '
        'Our state faces real challanges \u2014 from hurricaine recovery in the west to '
        'the mental health crisis effecting our communities. But i beleive that togather, '
        'we can build a safer, stronger North Carolna. Governer Stien has proposed a '
        '195 million dollar public saftey package that invests in law enforcment, '
        'violence preventation, and mental health treatement. Lets get it done.'
    )

    add_text_box(slide, Inches(0.6), Inches(1.85), Inches(3), Inches(0.4),
                 "Sample 1 \u2014 Short (quick demo)", font_size=13, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.4),
                 sample_text, font_size=14, color=RGBColor(0x44, 0x44, 0x44),
                 font_name="Consolas")

    # Sample 2 with speaker labels
    add_text_box(slide, Inches(0.6), Inches(4.85), Inches(4), Inches(0.4),
                 "Sample 2 \u2014 With speaker labels", font_size=13, color=DARK, bold=True)

    sample2_rect = add_rounded_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8),
                                    RGBColor(0xFD, 0xF6, 0xEC))
    sample2_rect.line.color.rgb = RGBColor(0xE8, 0xD5, 0xB0)
    sample2_rect.line.width = Pt(1)

    sample2_text = (
        'Gov. Stien: Thank you all for being hear today. The work were doing in western '
        'North Carolna after Hurricaine Helene has been nothing short of extrodinary.\n'
        'Reporter: Governer, can you give us an update on the road re-openings?\n'
        'Gov Stein: Absolutley. The Departement of Transporation has re-opened ninty eight '
        'percent of Helene-impacted state roads.'
    )

    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.6),
                 sample2_text, font_size=13, color=RGBColor(0x44, 0x44, 0x44),
                 font_name="Consolas")

    add_text_box(slide, Inches(0.4), Inches(7.0), Inches(12), Inches(0.4),
                 "NC DIT AI Office", font_size=9, color=RGBColor(0xBB, 0xBB, 0xBB))

    # ===== SLIDE 5: DEMO FLOW =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                 "Suggested Demo Flow", font_size=36, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    steps = [
        ("1", "Clips #1", "\"What clips came in this week about broadband?\"",
         "Basic search \u2014 gets the audience oriented", BLUE),
        ("2", "Remarks #1", "\"What has the Governor said about mental health and public safety?\"",
         "RAG synthesis with direct quotes and citations \u2014 the core value prop", GREEN),
        ("3", "Proofread", "Paste Sample 1 (short transcript with typos)",
         "Quick turnaround, every change explained with confidence levels", ORANGE),
        ("4", "Clips #3 (if time)", "\"Has the Governor said anything about rural internet access?\"",
         "The \"wow\" moment \u2014 semantic search finds \"broadband\" from \"rural internet\"", BLUE),
        ("5", "Remarks #3 (if time)", "\"How has the Governor talked about emergency preparedness?\"",
         "Cross-speech synthesis \u2014 pulls from August + November columns", GREEN),
    ]

    y = Inches(1.9)
    for num, label, prompt, note, color in steps:
        from pptx.enum.shapes import MSO_SHAPE
        # Step number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = num
        cp.font.size = Pt(18)
        cp.font.color.rgb = WHITE
        cp.font.bold = True
        cp.font.name = "Segoe UI"
        cp.alignment = PP_ALIGN.CENTER

        # Label
        add_text_box(slide, Inches(1.7), y - Inches(0.02), Inches(2), Inches(0.35),
                     label, font_size=18, color=WHITE, bold=True)

        # Prompt
        add_text_box(slide, Inches(3.8), y - Inches(0.02), Inches(8), Inches(0.35),
                     prompt, font_size=15, color=RGBColor(0xDD, 0xDD, 0xFF), italic=True)

        # Note
        add_text_box(slide, Inches(3.8), y + Inches(0.32), Inches(8), Inches(0.35),
                     note, font_size=11, color=RGBColor(0x99, 0xAA, 0xCC))

        y += Inches(1.05)

    add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
                 "NC DIT AI Office  |  March 2026", font_size=9,
                 color=RGBColor(0x66, 0x77, 0x99), alignment=PP_ALIGN.CENTER)

    # Save
    out = Path(__file__).parent / "demo-questions.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
