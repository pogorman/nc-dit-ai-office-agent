"""Generate presentation.pptx from the HTML slide deck content."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# -- Colors --
NAVY = RGBColor(0x0A, 0x24, 0x63)
BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xD8, 0x3B, 0x01)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
FAFA = RGBColor(0xFA, 0xFA, 0xFA)
DDD = RGBColor(0xDD, 0xDD, 0xDD)
PURPLE = RGBColor(0x7B, 0x2D, 0x8E)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE1)
AMBER_BORDER = RGBColor(0xF9, 0xA8, 0x25)
AMBER_TEXT = RGBColor(0xE6, 0x51, 0x00)
ASK_BG = RGBColor(0xF0, 0xF6, 0xFF)
FOOTER_GRAY = RGBColor(0xBB, 0xBB, 0xBB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def tb(slide, l, t, w, h, text, sz=18, color=DARK, bold=False, italic=False,
       align=PP_ALIGN.LEFT, font="Segoe UI", wrap=True):
    """Add a textbox and return (shape, text_frame)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    return box, tf


def rect(slide, l, t, w, h, fill, line_color=None, radius=0.05):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape


def circle(slide, l, t, size, fill, text, text_sz=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    ctf = shape.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.text = text
    cp.font.size = Pt(text_sz)
    cp.font.color.rgb = WHITE
    cp.font.bold = True
    cp.font.name = "Segoe UI"
    cp.alignment = PP_ALIGN.CENTER
    return shape


def footer(slide, left_text, right_text):
    tb(slide, Inches(0.6), Inches(7.05), Inches(5), Inches(0.3),
       left_text, sz=9, color=FOOTER_GRAY)
    tb(slide, Inches(7.5), Inches(7.05), Inches(5), Inches(0.3),
       right_text, sz=9, color=FOOTER_GRAY, align=PP_ALIGN.RIGHT)


def add_flow_row(slide, y, steps, accent_color):
    """Draw a 3-step flow row with boxes and arrows."""
    box_w = Inches(3.5)
    box_h = Inches(1.5)
    arrow_w = Inches(0.6)
    x = Inches(0.6)

    for i, (num, title, detail) in enumerate(steps):
        # Box
        r = rect(slide, x, y, box_w, box_h, FAFA, line_color=DDD)

        # Number circle
        circ_size = Inches(0.35)
        circle(slide, x + box_w / 2 - circ_size / 2, y + Inches(0.12), circ_size, accent_color, num, text_sz=13)

        # Title
        tb(slide, x + Inches(0.15), y + Inches(0.5), box_w - Inches(0.3), Inches(0.3),
           title, sz=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)

        # Detail
        tb(slide, x + Inches(0.15), y + Inches(0.8), box_w - Inches(0.3), Inches(0.6),
           detail, sz=11, color=GRAY, align=PP_ALIGN.CENTER)

        x += box_w

        # Arrow between boxes
        if i < len(steps) - 1:
            tb(slide, x, y + box_h / 2 - Inches(0.2), arrow_w, Inches(0.4),
               "\u2192", sz=26, color=FOOTER_GRAY, align=PP_ALIGN.CENTER)
            x += arrow_w


def add_insight_box(slide, y, bold_text, rest_text):
    """Add an amber insight callout box."""
    box_w = Inches(11.8)
    box_h = Inches(0.7)
    r = rect(slide, Inches(0.6), y, box_w, box_h, AMBER_BG, line_color=AMBER_BORDER)

    bx, tf = tb(slide, Inches(0.8), y + Inches(0.08), box_w - Inches(0.4), box_h - Inches(0.16),
                "", sz=13, color=DARK)
    p = tf.paragraphs[0]
    run_bold = p.add_run()
    run_bold.text = bold_text
    run_bold.font.size = Pt(13)
    run_bold.font.color.rgb = AMBER_TEXT
    run_bold.font.bold = True
    run_bold.font.name = "Segoe UI"
    run_rest = p.add_run()
    run_rest.text = rest_text
    run_rest.font.size = Pt(13)
    run_rest.font.color.rgb = DARK
    run_rest.font.name = "Segoe UI"


def add_whats_next_slide(prs, blank):
    """Slide 6: What's Next — Power Platform opportunities."""
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # Purple header bar
    cap_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    cap_bar.fill.solid()
    cap_bar.fill.fore_color.rgb = PURPLE
    cap_bar.line.fill.background()

    tb(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
       "What\u2019s Next \u2014 Power Platform", sz=32, color=WHITE, bold=True)
    tb(s, Inches(0.6), Inches(0.65), Inches(11), Inches(0.4),
       "Azure does the heavy lifting. Power Platform delivers the experience.",
       sz=15, color=RGBColor(0xE8, 0xD5, 0xF5))

    # Three opportunity cards
    cards = [
        ("Daily Digest Email",
         "Power Automate scheduled flow sends a morning clips summary via Outlook. "
         "Unblocks a stubbed feature with zero new Azure resources."),
        ("Dataverse for Metadata",
         "Move clips + remarks metadata to Dataverse. Staff can browse and edit "
         "in a model-driven app \u2014 no code needed."),
        ("Reliable Blob Trigger",
         "Power Automate \u201cWhen a blob is created\u201d trigger replaces the unreliable "
         "Flex Consumption blob trigger for remarks ingestion."),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.8)
    card_gap = Inches(0.5)
    total_w = card_w * 3 + card_gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (title, detail) in enumerate(cards):
        x = start_x + i * (card_w + card_gap)
        y = Inches(1.6)
        r = rect(s, x, y, card_w, card_h, FAFA, line_color=DDD)

        # Number circle
        circ_size = Inches(0.4)
        circle(s, x + card_w / 2 - circ_size / 2, y + Inches(0.2),
               circ_size, PURPLE, str(i + 1), text_sz=14)

        tb(s, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(0.3),
           title, sz=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(1.4),
           detail, sz=12, color=GRAY, align=PP_ALIGN.CENTER)

    # Insight box
    add_insight_box(s, Inches(4.8),
                    "Hybrid approach: ",
                    "AI Search, OpenAI, and the RAG pipeline stay in Azure where they belong. "
                    "Power Platform handles orchestration, notifications, and the data experience.")

    footer(s, "NC DIT AI Office", "6 / 9")


def add_architecture_slide(prs, blank):
    """Slide 7: Architecture overview — service topology."""
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
       "Architecture Overview", sz=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.6), Inches(0.75), Inches(12), Inches(0.3),
       "End-to-end service topology on Azure", sz=14,
       color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    BW, BH = Inches(1.95), Inches(1.15)
    AW = Inches(0.4)
    Y = Inches(1.7)
    FILL_A = RGBColor(0x15, 0x38, 0x70)
    FILL_B = RGBColor(0x1A, 0x40, 0x80)
    LINE_C = RGBColor(0x33, 0x55, 0x99)
    ARR_C = RGBColor(0x55, 0x77, 0xBB)

    cx = [Inches(0.35)]
    for _ in range(3):
        cx.append(cx[-1] + BW + AW)

    main_boxes = [
        ("Teams / Web", "Users ask questions\nor paste transcripts"),
        ("Copilot Studio", "Agent orchestration\nroutes to the right skill"),
        ("API Management", "Auth gateway\nsubscription key + routing"),
        ("Azure Functions", "Serverless compute\n6 TypeScript functions"),
    ]

    for i, (title, detail) in enumerate(main_boxes):
        fill = FILL_A if i % 2 == 0 else FILL_B
        rect(s, cx[i], Y, BW, BH, fill, line_color=LINE_C)
        tb(s, cx[i] + Inches(0.08), Y + Inches(0.15), BW - Inches(0.16), Inches(0.3),
           title, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, cx[i] + Inches(0.08), Y + Inches(0.5), BW - Inches(0.16), Inches(0.55),
           detail, sz=9, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
        if i < 3:
            tb(s, cx[i] + BW, Y + BH / 2 - Inches(0.15), AW, Inches(0.3),
               "\u2192", sz=22, color=ARR_C, align=PP_ALIGN.CENTER)

    # Arrow from Functions to backend services
    tb(s, cx[3] + BW, Y + BH / 2 - Inches(0.15), AW, Inches(0.3),
       "\u2192", sz=22, color=ARR_C, align=PP_ALIGN.CENTER)

    # Backend services (stacked vertically)
    SX = cx[3] + BW + AW
    SW, SH, SG = Inches(2.55), Inches(0.8), Inches(0.1)
    services = [
        ("Azure OpenAI", "GPT-4o  \u00b7  text-embedding-3-large", GREEN),
        ("AI Search", "Hybrid vector + keyword  \u00b7  semantic ranker", BLUE),
        ("Cosmos DB", "Serverless  \u00b7  metadata + full text", ORANGE),
    ]
    sy = Y - Inches(0.15)
    for title, detail, accent in services:
        rect(s, SX, sy, SW, SH, accent)
        tb(s, SX + Inches(0.08), sy + Inches(0.1), SW - Inches(0.16), Inches(0.28),
           title, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, SX + Inches(0.08), sy + Inches(0.4), SW - Inches(0.16), Inches(0.35),
           detail, sz=9, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
        sy += SH + SG

    # Supporting services row
    SPY = Inches(4.3)
    tb(s, Inches(0.6), SPY - Inches(0.35), Inches(12), Inches(0.3),
       "SUPPORTING SERVICES", sz=10, color=RGBColor(0x77, 0x88, 0xAA),
       bold=True, align=PP_ALIGN.CENTER)

    support = [
        ("Key Vault", "Function host key\nfor APIM only"),
        ("Blob Storage", "Remarks uploads\ntrigger ingestion"),
        ("VNet + Private Endpoints", "Network isolation\nCosmos + Blob"),
        ("Managed Identity", "DefaultAzureCredential\nzero secrets in code"),
    ]
    SPW, SPH, SPG = Inches(2.7), Inches(0.9), Inches(0.27)
    total = len(support) * SPW + (len(support) - 1) * SPG
    sx = (SLIDE_W - total) / 2
    SP_FILL = RGBColor(0x12, 0x2D, 0x55)
    SP_LINE = RGBColor(0x2A, 0x44, 0x77)
    for title, detail in support:
        rect(s, sx, SPY, SPW, SPH, SP_FILL, line_color=SP_LINE)
        tb(s, sx + Inches(0.08), SPY + Inches(0.1), SPW - Inches(0.16), Inches(0.25),
           title, sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tb(s, sx + Inches(0.08), SPY + Inches(0.38), SPW - Inches(0.16), Inches(0.45),
           detail, sz=9, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
        sx += SPW + SPG

    # Cost callout
    CY = Inches(5.6)
    rect(s, Inches(2.8), CY, Inches(7.7), Inches(0.6),
         RGBColor(0x0D, 0x22, 0x44), line_color=RGBColor(0x33, 0x55, 0x88))
    tb(s, Inches(3.0), CY + Inches(0.08), Inches(7.3), Inches(0.4),
       "~$120\u2013195/mo total  \u00b7  Serverless \u2014 scales to zero  \u00b7  100% consumption tier",
       sz=13, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    FTR = RGBColor(0x66, 0x77, 0x99)
    tb(s, Inches(0.6), Inches(7.05), Inches(5), Inches(0.3),
       "NC DIT AI Office", sz=9, color=FTR)
    tb(s, Inches(7.5), Inches(7.05), Inches(5), Inches(0.3),
       "7 / 9", sz=9, color=FTR, align=PP_ALIGN.RIGHT)


def add_request_flow_slide(prs, blank):
    """Slide 7: Request flow — one swim lane per capability."""
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    tb(s, Inches(0.6), Inches(0.3), Inches(12), Inches(0.55),
       "Request Flow", sz=34, color=NAVY, bold=True)
    tb(s, Inches(0.6), Inches(0.8), Inches(12), Inches(0.3),
       "What happens when a user sends a message", sz=14, color=LIGHT_GRAY)

    BW, BH = Inches(2.35), Inches(1.2)
    AW = Inches(0.3)
    START_X = Inches(1.4)

    lanes = [
        ("News\nClips", BLUE, Inches(1.4), [
            ("User asks in Teams", "\"What clips came in\nabout broadband?\""),
            ("APIM routes request", "Subscription key auth\n\u2192 clips-query function"),
            ("AI Search", "Hybrid vector + keyword\nsemantic reranker"),
            ("Return matching clips", "Title, date, outlet,\nGovernor mention"),
        ]),
        ("Remarks\nSearch", GREEN, Inches(3.0), [
            ("User asks in Teams", "\"What has the Governor\nsaid about teacher pay?\""),
            ("APIM routes request", "Subscription key auth\n\u2192 remarks-query function"),
            ("AI Search + GPT-4o", "Retrieve chunks, then\nsynthesize with quotes"),
            ("Return summary", "Written synthesis with\ndirect quotes + sources"),
        ]),
        ("Proof-\nread", ORANGE, Inches(4.6), [
            ("User pastes transcript", "Raw text with\nASR/OCR errors"),
            ("APIM routes request", "Subscription key auth\n\u2192 proofread function"),
            ("GPT-4o proofreads", "Fix errors, preserve\nspeaking style"),
            ("Return corrections", "Clean text + change list\nwith confidence levels"),
        ]),
    ]

    for label, color, y, steps in lanes:
        # Color bar on left
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.4), y, Inches(0.07), BH)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Lane label
        tb(s, Inches(0.52), y + Inches(0.3), Inches(0.85), Inches(0.6),
           label, sz=12, color=color, bold=True, align=PP_ALIGN.CENTER)

        # Step boxes
        x = START_X
        for i, (title, detail) in enumerate(steps):
            rect(s, x, y, BW, BH, FAFA, line_color=DDD)
            circle(s, x + Inches(0.08), y + Inches(0.08), Inches(0.26),
                   color, str(i + 1), text_sz=10)
            tb(s, x + Inches(0.4), y + Inches(0.08), BW - Inches(0.5), Inches(0.4),
               title, sz=11, color=DARK, bold=True)
            tb(s, x + Inches(0.4), y + Inches(0.5), BW - Inches(0.5), Inches(0.6),
               detail, sz=9, color=GRAY)
            x += BW
            if i < len(steps) - 1:
                tb(s, x, y + BH / 2 - Inches(0.12), AW, Inches(0.25),
                   "\u2192", sz=18, color=FOOTER_GRAY, align=PP_ALIGN.CENTER)
                x += AW

    footer(s, "NC DIT AI Office", "8 / 9")


def add_data_flow_slide(prs, blank):
    """Slide 8: Data flow — how data enters and exits the system."""
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    tb(s, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
       "Data Flow", sz=34, color=NAVY, bold=True)
    tb(s, Inches(0.6), Inches(0.75), Inches(12), Inches(0.3),
       "How data enters the system and how users get answers", sz=14, color=LIGHT_GRAY)

    def section_line(y, text):
        tb(s, Inches(0.6), y, Inches(3), Inches(0.3), text, sz=12, color=NAVY, bold=True)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), y + Inches(0.28), Inches(12.1), Inches(0.02))
        bar.fill.solid()
        bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()

    BW, BH = Inches(2.35), Inches(0.75)
    AW = Inches(0.28)

    def flow_row(y, label, label_color, steps, accent_first=True):
        tb(s, Inches(0.45), y + Inches(0.15), Inches(0.85), Inches(0.45),
           label, sz=10, color=label_color, bold=True, align=PP_ALIGN.CENTER)
        x = Inches(1.4)
        for i, (title, detail) in enumerate(steps):
            is_accent = (accent_first and i == 0) or (not accent_first and i == len(steps) - 1)
            if is_accent:
                rect(s, x, y, BW, BH, label_color)
                tb(s, x + Inches(0.08), y + Inches(0.08), BW - Inches(0.16), Inches(0.28),
                   title, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
                tb(s, x + Inches(0.08), y + Inches(0.38), BW - Inches(0.16), Inches(0.32),
                   detail, sz=8, color=RGBColor(0xEE, 0xEE, 0xEE), align=PP_ALIGN.CENTER)
            else:
                rect(s, x, y, BW, BH, FAFA, line_color=DDD)
                tb(s, x + Inches(0.08), y + Inches(0.08), BW - Inches(0.16), Inches(0.28),
                   title, sz=10, color=DARK, bold=True, align=PP_ALIGN.CENTER)
                tb(s, x + Inches(0.08), y + Inches(0.38), BW - Inches(0.16), Inches(0.32),
                   detail, sz=8, color=GRAY, align=PP_ALIGN.CENTER)
            x += BW
            if i < len(steps) - 1:
                tb(s, x, y + BH / 2 - Inches(0.1), AW, Inches(0.2),
                   "\u2192", sz=16, color=FOOTER_GRAY, align=PP_ALIGN.CENTER)
                x += AW

    # DATA IN
    section_line(Inches(1.15), "DATA IN \u2014 Ingestion")

    flow_row(Inches(1.55), "CLIPS", BLUE, [
        ("Timer \u2014 7 AM daily", "Azure Function\nclips-ingest"),
        ("Scrape + Web Search", "governor.nc.gov +\nBing grounding (20+ outlets)"),
        ("Embed", "Azure OpenAI\ntext-embedding-3-large"),
        ("Store + Index", "Cosmos DB (metadata)\nAI Search (vectors)"),
    ])

    flow_row(Inches(2.5), "REMARKS", GREEN, [
        ("Blob Upload trigger", "Azure Function\nremarks-ingest"),
        ("Chunk into paragraphs", "Split speech into\nparagraph-sized chunks"),
        ("Embed", "Azure OpenAI\ntext-embedding-3-large"),
        ("Store + Index", "Cosmos DB (chunks)\nAI Search (vectors)"),
    ])

    # DATA OUT
    section_line(Inches(3.55), "DATA OUT \u2014 Queries")

    flow_row(Inches(3.95), "CLIPS", BLUE, [
        ("User asks question", "Copilot Studio \u2192 APIM\n\u2192 clips-query"),
        ("AI Search", "Hybrid vector + keyword\nsemantic reranker"),
        ("Return results", "Matching clips with\ntitle, date, outlet"),
        ("Display in Teams", "Agent formats clips\nfor the user"),
    ], accent_first=False)

    flow_row(Inches(4.9), "REMARKS", GREEN, [
        ("User asks question", "Copilot Studio \u2192 APIM\n\u2192 remarks-query"),
        ("AI Search", "Retrieve top matching\nchunks from speeches"),
        ("GPT-4o synthesizes", "Write summary with\ndirect quotes + sources"),
        ("Display in Teams", "Agent shows synthesis\nwith citations"),
    ], accent_first=False)

    flow_row(Inches(5.85), "PROOFREAD", ORANGE, [
        ("User pastes text", "Copilot Studio \u2192 APIM\n\u2192 proofread"),
        ("GPT-4o proofreads", "Fix errors, preserve\nspeaking style"),
        ("Return corrections", "Clean text + change list\nwith confidence levels"),
        ("Display in Teams", "Agent shows corrected\ntext + all changes"),
    ], accent_first=False)

    footer(s, "NC DIT AI Office", "9 / 9")


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ================================================================
    # SLIDE 1 — TITLE
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    tb(s, Inches(1), Inches(1.8), Inches(11.3), Inches(1.4),
       "NC Governor's\nCommunications Agent", sz=48, color=WHITE, bold=True,
       align=PP_ALIGN.CENTER)

    tb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8),
       "AI-powered news monitoring, remarks search, and transcript proofreading\nfor the Governor's Communications Office",
       sz=22, color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    tb(s, Inches(1), Inches(4.7), Inches(11.3), Inches(0.4),
       "NC DIT AI Office  \u00b7  Microsoft Azure  \u00b7  Copilot Studio",
       sz=14, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

    tb(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.4),
       "March 2026", sz=14, color=RGBColor(0x77, 0x88, 0xAA), align=PP_ALIGN.CENTER)

    # ================================================================
    # SLIDE 2 — THE ASK
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    tb(s, Inches(0.6), Inches(0.4), Inches(10), Inches(0.7),
       "What the Comms Office Asked For", sz=34, color=NAVY, bold=True)

    # Request 1 card — News Clips
    r1 = rect(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.5), ASK_BG)
    bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(0.08), Inches(1.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = NAVY
    bar1.line.fill.background()

    tb(s, Inches(0.95), Inches(1.4), Inches(11), Inches(0.3),
       "REQUEST 1 \u2014 NEWS CLIPS", sz=11, color=NAVY, bold=True)
    tb(s, Inches(0.95), Inches(1.7), Inches(11.5), Inches(0.9),
       '\u201cHow could we automate the process of identifying mentions of Governor Stein in the news and collecting the outlet, title, first paragraph, and first mention of Governor Stein in the article?\u201d',
       sz=16, color=RGBColor(0x33, 0x33, 0x33), italic=True)

    # Request 2 card — Proofreading
    r2 = rect(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.1), ASK_BG)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.0), Inches(0.08), Inches(1.1))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = NAVY
    bar2.line.fill.background()

    tb(s, Inches(0.95), Inches(3.1), Inches(11), Inches(0.3),
       "REQUEST 2 \u2014 PROOFREADING", sz=11, color=NAVY, bold=True)
    tb(s, Inches(0.95), Inches(3.4), Inches(11.5), Inches(0.6),
       '\u201cIs there a way to use AI for proofreading of faulty transcripts?\u201d',
       sz=16, color=RGBColor(0x33, 0x33, 0x33), italic=True)

    # Request 3 card — Remarks
    r3 = rect(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(1.5), ASK_BG)
    bar3 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.3), Inches(0.08), Inches(1.5))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = NAVY
    bar3.line.fill.background()

    tb(s, Inches(0.95), Inches(4.4), Inches(11), Inches(0.3),
       "REQUEST 3 \u2014 REMARKS", sz=11, color=NAVY, bold=True)
    tb(s, Inches(0.95), Inches(4.7), Inches(11.5), Inches(0.9),
       '\u201cHow could we create a useful search + retrieval function for existing language on a given topic? For example, what is the language we\u2019ve used to talk about clean tech across a variety of remarks?\u201d',
       sz=16, color=RGBColor(0x33, 0x33, 0x33), italic=True)

    # Built badge
    badge = rect(s, Inches(0.6), Inches(6.2), Inches(7.5), Inches(0.7), GREEN)
    _, btf = tb(s, Inches(0.8), Inches(6.25), Inches(7.1), Inches(0.6),
                "\u2713  All three capabilities \u2014 built, deployed, and live in Teams",
                sz=18, color=WHITE, bold=True)

    footer(s, "NC DIT AI Office", "2 / 9")

    # ================================================================
    # SLIDE 3 — NEWS CLIPS
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # Blue header bar
    cap_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    cap_bar.fill.solid()
    cap_bar.fill.fore_color.rgb = BLUE
    cap_bar.line.fill.background()

    tb(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
       "News Clips", sz=32, color=WHITE, bold=True)
    tb(s, Inches(0.6), Inches(0.65), Inches(11), Inches(0.4),
       "Monitors 20+ news outlets and Governor press releases, then lets staff search or browse them instantly.",
       sz=15, color=RGBColor(0xDD, 0xEE, 0xFF))

    # Ingestion flow
    tb(s, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
       "HOW DATA GETS IN \u2014 RUNS DAILY AT 7 AM", sz=10, color=LIGHT_GRAY, bold=True)

    add_flow_row(s, Inches(1.85), [
        ("1", "Scrape + Search", "Azure Function scrapes governor.nc.gov AND searches the web via Azure OpenAI + Bing grounding (WRAL, AP, 20+ outlets)"),
        ("2", "Embed", "Azure OpenAI converts each article into a vector embedding \u2014 a numeric fingerprint of its meaning"),
        ("3", "Store + Index", "Article + embedding saved in Cosmos DB (browsing) and Azure AI Search (hybrid search)"),
    ], BLUE)

    # Query flow
    tb(s, Inches(0.6), Inches(3.6), Inches(10), Inches(0.3),
       'HOW STAFF ASK \u2014 "WHAT\'S IN THE NEWS ABOUT BROADBAND?"', sz=10, color=LIGHT_GRAY, bold=True)

    add_flow_row(s, Inches(3.95), [
        ("1", "Ask in Teams", "Staff types a question in the Copilot Studio chatbot"),
        ("2", "Hybrid Search", "AI Search matches by meaning and keywords, then re-ranks with semantic ranker"),
        ("3", "Results", "Matching clips returned with title, date, outlet, and the Governor mention"),
    ], BLUE)

    # Insight
    add_insight_box(s, Inches(5.75),
                    "Key: ",
                    '"police pay" finds an article titled "Governor Stein Calls for Law Enforcement Raises" \u2014 because the meaning matches, not just the words.')

    footer(s, "NC DIT AI Office", "3 / 9")

    # ================================================================
    # SLIDE 4 — REMARKS SEARCH
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # Green header bar
    cap_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    cap_bar.fill.solid()
    cap_bar.fill.fore_color.rgb = GREEN
    cap_bar.line.fill.background()

    tb(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
       "Remarks Search", sz=32, color=WHITE, bold=True)
    tb(s, Inches(0.6), Inches(0.65), Inches(11), Inches(0.4),
       "Search the Governor's speeches by topic. AI reads all matching passages and writes a summary with direct quotes.",
       sz=15, color=RGBColor(0xDD, 0xFF, 0xDD))

    # Ingestion flow
    tb(s, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
       "HOW SPEECHES GET IN \u2014 UPLOAD TRIGGERS PROCESSING", sz=10, color=LIGHT_GRAY, bold=True)

    add_flow_row(s, Inches(1.85), [
        ("1", "Upload", "Drop a speech file into Blob Storage \u2014 triggers an Azure Function automatically"),
        ("2", "Chunk + Embed", "Function splits into paragraphs, embeds each chunk, tags with date/event/venue"),
        ("3", "Index", "Chunks stored in Cosmos DB + Azure AI Search for retrieval"),
    ], GREEN)

    # Query flow
    tb(s, Inches(0.6), Inches(3.6), Inches(10), Inches(0.3),
       'HOW STAFF ASK \u2014 "WHAT HAS THE GOVERNOR SAID ABOUT TEACHER PAY?"', sz=10, color=LIGHT_GRAY, bold=True)

    add_flow_row(s, Inches(3.95), [
        ("1", "Ask in Teams", "Staff types a topic question in the Copilot Studio chatbot"),
        ("2", "Hybrid Search", "Finds top matching speech chunks across all remarks via vector + keyword search"),
        ("3", "GPT-4o Synthesizes", "AI reads the chunks and writes a synthesis with direct quotes, dates, and sources"),
    ], GREEN)

    # Insight
    add_insight_box(s, Inches(5.75),
                    "This is the RAG pattern ",
                    "(Retrieval-Augmented Generation): we retrieve real source material first, then ask AI to summarize only what it found. The AI never invents quotes \u2014 it cites real text from real speeches.")

    footer(s, "NC DIT AI Office", "4 / 9")

    # ================================================================
    # SLIDE 5 — PROOFREAD + PLATFORM
    # ================================================================
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # Orange header bar
    cap_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.2))
    cap_bar.fill.solid()
    cap_bar.fill.fore_color.rgb = ORANGE
    cap_bar.line.fill.background()

    tb(s, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
       "Transcript Proofreading", sz=32, color=WHITE, bold=True)
    tb(s, Inches(0.6), Inches(0.65), Inches(11), Inches(0.4),
       "Paste a rough transcript, get back a clean version with every change explained.",
       sz=15, color=RGBColor(0xFF, 0xDD, 0xCC))

    # Flow
    tb(s, Inches(0.6), Inches(1.45), Inches(10), Inches(0.3),
       "HOW IT WORKS \u2014 ONE STEP, NO DATA STORED", sz=10, color=LIGHT_GRAY, bold=True)

    add_flow_row(s, Inches(1.85), [
        ("1", "Paste Text", "Staff pastes a raw transcript from voice-to-text or a scanned document into the chatbot"),
        ("2", "GPT-4o Proofreads", "Fixes ASR/OCR errors, punctuation, proper nouns \u2014 never changes meaning or speaking style"),
        ("3", "Structured Output", "Returns corrected text + a list of every change with a reason and confidence level"),
    ], ORANGE)

    # Insight
    add_insight_box(s, Inches(3.6),
                    "Key: ",
                    "Like handing a rough draft to the best copy editor in the building \u2014 except they work instantly and show you every red-pen mark with an explanation.")

    # Three bottom cards
    cards = [
        ("Nothing Stored", "Text in, clean text out.\nNo data touches the database."),
        ("~$120\u2013195/mo", "Total platform cost.\nServerless \u2014 scales to zero when idle."),
        ("Zero Passwords in Code", "Managed identity everywhere.\nData stays in the state's Azure tenant."),
    ]

    card_w = Inches(3.6)
    card_h = Inches(1.4)
    card_gap = Inches(0.5)
    total_w = card_w * 3 + card_gap * 2
    start_x = (SLIDE_W - total_w) / 2

    for i, (title, detail) in enumerate(cards):
        x = start_x + i * (card_w + card_gap)
        y = Inches(4.7)
        r = rect(s, x, y, card_w, card_h, FAFA, line_color=DDD)

        tb(s, x + Inches(0.2), y + Inches(0.25), card_w - Inches(0.4), Inches(0.3),
           title, sz=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        tb(s, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.7),
           detail, sz=11, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s, "NC DIT AI Office", "5 / 9")

    # ================================================================
    # SLIDE 6 — WHAT'S NEXT — POWER PLATFORM
    # ================================================================
    add_whats_next_slide(prs, blank)

    # ================================================================
    # SLIDE 7 — ARCHITECTURE OVERVIEW
    # ================================================================
    add_architecture_slide(prs, blank)

    # ================================================================
    # SLIDE 8 — REQUEST FLOW
    # ================================================================
    add_request_flow_slide(prs, blank)

    # ================================================================
    # SLIDE 9 — DATA FLOW
    # ================================================================
    add_data_flow_slide(prs, blank)

    # Save
    out = Path(__file__).parent / "html" / "presentation.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
